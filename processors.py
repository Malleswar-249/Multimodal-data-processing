"""
File processors for different file types
"""
import os
import base64
from typing import Dict, List, Tuple
import PyPDF2
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import speech_recognition as sr
from pydub import AudioSegment
from youtube_transcript_api import YouTubeTranscriptApi
import yt_dlp
import io


class TextProcessor:
    """Process text-based files"""
    
    @staticmethod
    def process_pdf(file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            return f"Error processing PDF: {str(e)}"
    
    @staticmethod
    def process_docx(file_path: str) -> str:
        """Extract text from DOCX file"""
        try:
            doc = Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text.strip()
        except Exception as e:
            return f"Error processing DOCX: {str(e)}"
    
    @staticmethod
    def process_pptx(file_path: str) -> str:
        """Extract text from PPTX file"""
        try:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text).strip()
        except Exception as e:
            return f"Error processing PPTX: {str(e)}"
    
    @staticmethod
    def process_txt(file_path: str) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read().strip()
        except Exception as e:
            return f"Error processing TXT: {str(e)}"
    
    @staticmethod
    def process_md(file_path: str) -> str:
        """Extract text from Markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read().strip()
        except Exception as e:
            return f"Error processing MD: {str(e)}"


class ImageProcessor:
    """Process image files"""
    
    @staticmethod
    def process_image(file_path: str) -> Tuple[str, str]:
        """Extract text from image using OCR and return base64 encoded image"""
        try:
            # Load image
            image = Image.open(file_path)
            
            # Convert to base64 for storage
            with open(file_path, 'rb') as img_file:
                img_data = base64.b64encode(img_file.read()).decode('utf-8')
            
            # Extract text using OCR
            try:
                ocr_text = pytesseract.image_to_string(image)
            except Exception:
                ocr_text = "OCR text extraction failed. Image stored for visual reference."
            
            description = f"Image file: {os.path.basename(file_path)}\nExtracted text from image:\n{ocr_text}"
            
            return description.strip(), img_data
        except Exception as e:
            return f"Error processing image: {str(e)}", ""


class AudioVideoProcessor:
    """Process audio and video files"""
    
    @staticmethod
    def process_mp3(file_path: str) -> str:
        """Extract text from MP3 audio file using speech recognition"""
        try:
            # Convert MP3 to WAV for speech recognition
            audio = AudioSegment.from_mp3(file_path)
            wav_path = file_path.replace('.mp3', '.wav')
            audio.export(wav_path, format="wav")
            
            # Transcribe audio
            recognizer = sr.Recognizer()
            with sr.AudioFile(wav_path) as source:
                audio_data = recognizer.record(source)
            
            try:
                text = recognizer.recognize_google(audio_data)
            except sr.UnknownValueError:
                text = "Could not understand audio"
            except sr.RequestError:
                text = "Speech recognition service unavailable"
            
            # Clean up temporary WAV file
            if os.path.exists(wav_path):
                os.remove(wav_path)
            
            return f"Audio transcription:\n{text}"
        except Exception as e:
            return f"Error processing MP3: {str(e)}"
    
    @staticmethod
    def process_mp4(file_path: str) -> str:
        """Extract audio from MP4 and transcribe"""
        try:
            # Extract audio from video
            audio = AudioSegment.from_file(file_path, format="mp4")
            wav_path = file_path.replace('.mp4', '.wav')
            audio.export(wav_path, format="wav")
            
            # Transcribe audio
            recognizer = sr.Recognizer()
            with sr.AudioFile(wav_path) as source:
                audio_data = recognizer.record(source)
            
            try:
                text = recognizer.recognize_google(audio_data)
            except sr.UnknownValueError:
                text = "Could not understand audio from video"
            except sr.RequestError:
                text = "Speech recognition service unavailable"
            
            # Clean up temporary WAV file
            if os.path.exists(wav_path):
                os.remove(wav_path)
            
            return f"Video transcription:\n{text}"
        except Exception as e:
            return f"Error processing MP4: {str(e)}"
    
    @staticmethod
    def process_youtube(url: str) -> str:
        """Extract transcript from YouTube video"""
        try:
            # Extract video ID from URL
            video_id = url.split('v=')[-1].split('&')[0]
            
            # Try to get transcript
            try:
                transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
                transcript_text = " ".join([item['text'] for item in transcript_list])
                return f"YouTube video transcript:\n{transcript_text}"
            except Exception:
                # If transcript not available, try to download and process audio
                return f"YouTube video URL processed: {url}\n(Transcript not available, audio processing would be required)"
        except Exception as e:
            return f"Error processing YouTube URL: {str(e)}"


class FileProcessor:
    """Main file processor that routes to appropriate handlers"""
    
    def __init__(self):
        self.text_processor = TextProcessor()
        self.image_processor = ImageProcessor()
        self.audio_video_processor = AudioVideoProcessor()
    
    def process_file(self, file_path: str = None, file_data=None, file_name: str = None, file_type: str = None) -> Dict[str, any]:
        """
        Process file based on its type
        Returns dict with 'text', 'metadata', and 'image_data' (for images)
        """
        if file_path:
            file_name = os.path.basename(file_path)
            file_ext = os.path.splitext(file_name)[1].lower()
        elif file_data and file_name:
            file_ext = os.path.splitext(file_name)[1].lower()
        elif file_type:
            file_ext = f".{file_type.lower()}"
        else:
            return {"text": "", "metadata": {}, "image_data": ""}
        
        metadata = {
            "file_name": file_name,
            "file_type": file_ext,
            "file_path": file_path if file_path else "uploaded_file"
        }
        
        result = {"text": "", "metadata": metadata, "image_data": ""}
        
        try:
            if file_ext in ['.pdf']:
                if file_path:
                    result["text"] = self.text_processor.process_pdf(file_path)
                else:
                    # For uploaded files, save temporarily
                    temp_path = f"temp_{file_name}"
                    with open(temp_path, 'wb') as f:
                        f.write(file_data.getvalue())
                    result["text"] = self.text_processor.process_pdf(temp_path)
                    os.remove(temp_path)
                    
            elif file_ext in ['.docx']:
                if file_path:
                    result["text"] = self.text_processor.process_docx(file_path)
                else:
                    temp_path = f"temp_{file_name}"
                    with open(temp_path, 'wb') as f:
                        f.write(file_data.getvalue())
                    result["text"] = self.text_processor.process_docx(temp_path)
                    os.remove(temp_path)
                    
            elif file_ext in ['.pptx']:
                if file_path:
                    result["text"] = self.text_processor.process_pptx(file_path)
                else:
                    temp_path = f"temp_{file_name}"
                    with open(temp_path, 'wb') as f:
                        f.write(file_data.getvalue())
                    result["text"] = self.text_processor.process_pptx(temp_path)
                    os.remove(temp_path)
                    
            elif file_ext in ['.txt']:
                if file_path:
                    result["text"] = self.text_processor.process_txt(file_path)
                else:
                    result["text"] = file_data.read().decode('utf-8')
                    
            elif file_ext in ['.md']:
                if file_path:
                    result["text"] = self.text_processor.process_md(file_path)
                else:
                    result["text"] = file_data.read().decode('utf-8')
                    
            elif file_ext in ['.png', '.jpg', '.jpeg']:
                if file_path:
                    text, img_data = self.image_processor.process_image(file_path)
                    result["text"] = text
                    result["image_data"] = img_data
                else:
                    temp_path = f"temp_{file_name}"
                    with open(temp_path, 'wb') as f:
                        f.write(file_data.getvalue())
                    text, img_data = self.image_processor.process_image(temp_path)
                    result["text"] = text
                    result["image_data"] = img_data
                    os.remove(temp_path)
                    
            elif file_ext in ['.mp3']:
                if file_path:
                    result["text"] = self.audio_video_processor.process_mp3(file_path)
                else:
                    temp_path = f"temp_{file_name}"
                    with open(temp_path, 'wb') as f:
                        f.write(file_data.getvalue())
                    result["text"] = self.audio_video_processor.process_mp3(temp_path)
                    os.remove(temp_path)
                    
            elif file_ext in ['.mp4']:
                if file_path:
                    result["text"] = self.audio_video_processor.process_mp4(file_path)
                else:
                    temp_path = f"temp_{file_name}"
                    with open(temp_path, 'wb') as f:
                        f.write(file_data.getvalue())
                    result["text"] = self.audio_video_processor.process_mp4(temp_path)
                    os.remove(temp_path)
                    
            else:
                result["text"] = f"Unsupported file type: {file_ext}"
                
        except Exception as e:
            result["text"] = f"Error processing file: {str(e)}"
        
        return result
    
    def process_youtube_url(self, url: str) -> Dict[str, any]:
        """Process YouTube URL"""
        metadata = {
            "file_name": url,
            "file_type": ".youtube",
            "file_path": url
        }
        
        return {
            "text": self.audio_video_processor.process_youtube(url),
            "metadata": metadata,
            "image_data": ""
        }

